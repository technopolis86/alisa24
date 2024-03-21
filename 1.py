from pptx import Presentation
from random import sample
import random

prs = Presentation()

lst = [i for i in dir(random) if '_' not in i][6:]
sp = sample(lst, 5)
print(sp)
title_slide_layout = prs.slide_layouts[0]
# создаем заголовочный слайд
slide = prs.slides.add_slide(title_slide_layout)
# создаем у слайда заголовок и текст
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Библиотека Random"
for i in [random.randint, random.sample]:
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[0]
    title.text = i.__name__
    subtitle = slide.placeholders[1]
    subtitle.text = i.__doc__

prs.save('test.pptx')